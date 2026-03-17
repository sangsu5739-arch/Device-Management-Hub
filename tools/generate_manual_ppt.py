# -*- coding: utf-8 -*-
"""
Universal Device Studio (UDS) - User Manual PPT Generator  (v2)
Run: python tools/generate_manual_ppt.py
Output: docs/UDS_UserManual.pptx
20 slides + image placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, re

# ── Color Palette ─────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x1A, 0x1F, 0x2E)
C_PANEL   = RGBColor(0x24, 0x28, 0x38)
C_PANEL2  = RGBColor(0x2B, 0x2F, 0x42)
C_ACCENT  = RGBColor(0x4A, 0x90, 0xE2)
C_CYAN    = RGBColor(0x00, 0xD2, 0xFF)
C_ORANGE  = RGBColor(0xFF, 0xA5, 0x00)
C_GREEN   = RGBColor(0x3C, 0xC6, 0x8A)
C_PURPLE  = RGBColor(0xAA, 0x80, 0xFF)
C_YELLOW  = RGBColor(0xF0, 0xC0, 0x40)
C_RED     = RGBColor(0xE0, 0x50, 0x50)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT    = RGBColor(0xE0, 0xE0, 0xE0)
C_SUBTEXT = RGBColor(0xA0, 0xB0, 0xC0)
C_BORDER  = RGBColor(0x3A, 0x40, 0x58)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

MODULE_COLORS = {
    "FTDI Verifier":  C_ACCENT,
    "INA228":         C_CYAN,
    "PI6CG18201":     C_GREEN,
    "ADS1018":        C_ORANGE,
}


# ── Primitives ────────────────────────────────────────────────────────────────

def _clean(text):
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '',
                  str(text).replace('\n', ' '))


def set_bg(slide, color=None):
    color = color or C_BG
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def rect(slide, l, t, w, h, fill, line=None, lw=Pt(1), radius=False):
    shp = slide.shapes.add_shape(5 if radius else 1, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line: shp.line.color.rgb = line; shp.line.width = lw
    else: shp.line.fill.background()
    return shp


def tb(slide, text, l, t, w, h, sz=12, bold=False, color=C_TEXT,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = _clean(text)
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return box


def title_bar(slide, title, subtitle='', accent=C_ACCENT):
    rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_PANEL)
    rect(slide, 0, 0, Inches(0.1), Inches(1.1), accent)
    tb(slide, title, Inches(0.2), Inches(0.08), Inches(11), Inches(0.62),
       sz=26, bold=True, color=C_WHITE)
    if subtitle:
        tb(slide, subtitle, Inches(0.2), Inches(0.68), Inches(11), Inches(0.36),
           sz=12, color=C_SUBTEXT)
    # Slide branding right
    tb(slide, "Universal Device Studio  v1.0.0",
       Inches(10.5), Inches(0.08), Inches(2.7), Inches(0.35),
       sz=9, color=C_BORDER, align=PP_ALIGN.RIGHT)


def img_placeholder(slide, l, t, w, h, label='[Image]', color=C_PANEL2):
    rect(slide, l, t, w, h, color, C_BORDER, Pt(1.2))
    # Diagonal cross lines (simulate image box)
    tb(slide, label, l + Inches(0.1), t + h/2 - Pt(10),
       w - Inches(0.2), Inches(0.5),
       sz=9, color=C_SUBTEXT, align=PP_ALIGN.CENTER, italic=True)


def card(slide, l, t, w, h, title, body_lines, accent=C_ACCENT,
         title_sz=12, body_sz=10.5):
    rect(slide, l, t, w, h, C_PANEL, accent, Pt(1))
    rect(slide, l, t, Inches(0.07), h, accent)
    tb(slide, title, l+Inches(0.12), t+Inches(0.06), w-Inches(0.2), Inches(0.35),
       sz=title_sz, bold=True, color=accent)
    rect(slide, l+Inches(0.07), t+Inches(0.42), w-Inches(0.1), Pt(0.8), accent)
    y = t + Inches(0.48)
    lh = Inches(0.3)
    for line in body_lines:
        tb(slide, '  > '+line if not line.startswith('[') else line,
           l+Inches(0.12), y, w-Inches(0.2), lh, sz=body_sz, color=C_TEXT)
        y += lh


def badge(slide, l, t, w, h, text, bg=C_ACCENT, fg=C_WHITE, sz=11, bold=True):
    rect(slide, l, t, w, h, bg, radius=True)
    tb(slide, text, l, t, w, h, sz=sz, bold=bold, color=fg,
       align=PP_ALIGN.CENTER)


def section_header(slide, text, y=Inches(1.2), color=C_ACCENT):
    tb(slide, text, Inches(0.3), y, Inches(12.5), Inches(0.38),
       sz=13, bold=True, color=color)
    rect(slide, Inches(0.3), y+Inches(0.38), Inches(12.5), Pt(1), color)


def bullet_list(slide, items, l, t, w, sz=11, color=C_TEXT, gap=Inches(0.33)):
    y = t
    for item in items:
        tb(slide, '  > '+item, l, y, w, gap, sz=sz, color=color)
        y += gap
    return y


def two_col_header(slide, left_title, right_title, lc=C_ACCENT, rc=C_CYAN):
    mid = Inches(6.55)
    rect(slide, Inches(0.3), Inches(1.2), Inches(6.1), Inches(0.38), lc)
    tb(slide, left_title, Inches(0.3), Inches(1.2), Inches(6.1), Inches(0.38),
       sz=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    rect(slide, mid, Inches(1.2), Inches(6.6), Inches(0.38), rc)
    tb(slide, right_title, mid, Inches(1.2), Inches(6.6), Inches(0.38),
       sz=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)


# ── Slides ───────────────────────────────────────────────────────────────────

def s01_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    # Accent panel
    rect(s, 0, Inches(2.6), SLIDE_W, Inches(2.4), C_PANEL)
    rect(s, 0, Inches(2.6), Inches(0.15), Inches(2.4), C_ACCENT)

    tb(s, "Universal Device Studio", Inches(0.4), Inches(2.75),
       Inches(8.5), Inches(0.95), sz=42, bold=True, color=C_WHITE)
    tb(s, "FT232H MPSSE 기반 멀티 프로토콜 하드웨어 검증 플랫폼",
       Inches(0.4), Inches(3.72), Inches(8.5), Inches(0.45),
       sz=16, color=C_CYAN)
    tb(s, "I2C / SPI / JTAG / UART / GPIO  —  All-in-One",
       Inches(0.4), Inches(4.22), Inches(8.5), Inches(0.38),
       sz=13, color=C_SUBTEXT)

    badge(s, Inches(0.4), Inches(4.75), Inches(1.3), Inches(0.38),
          "v1.0.0", C_ACCENT)
    badge(s, Inches(1.85), Inches(4.75), Inches(1.8), Inches(0.38),
          "STATSChipPAC", C_PANEL2, C_SUBTEXT, sz=10, bold=False)

    img_placeholder(s, Inches(9.1), Inches(1.0), Inches(4.0), Inches(5.5),
                    '[스크린샷: 다크 테마 메인 창 전체\n(연결된 상태, FTDI Verifier 탭 활성)]')
    img_placeholder(s, Inches(0.4), Inches(6.2), Inches(1.2), Inches(1.0),
                    '[UDS 아이콘]')

    tb(s, "STATSChipPAC  |  Hardware Verification & Control Platform",
       Inches(1.8), Inches(6.95), Inches(11), Inches(0.38),
       sz=10, color=C_BORDER)


def s02_toc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "목차", "Table of Contents")

    sections = [
        ("1", "제품 소개 및 시스템 요구사항",  "Slide 3–4",  C_ACCENT),
        ("2", "설치 및 초기 설정",             "Slide 5",    C_CYAN),
        ("3", "메인 인터페이스 개요",           "Slide 6–7",  C_GREEN),
        ("4", "FTDI Verifier 모듈",            "Slide 8–12", C_ACCENT),
        ("5", "INA228 전력 모니터",             "Slide 13–15",C_CYAN),
        ("6", "PI6CG18201 클럭 제너레이터",     "Slide 16–17",C_GREEN),
        ("7", "ADS1018 ADC 모니터",            "Slide 18–19",C_ORANGE),
        ("8", "트러블슈팅 및 주의사항",         "Slide 20",   C_RED),
    ]
    cx, cy, cw, ch = Inches(0.3), Inches(1.3), Inches(5.8), Inches(0.68)
    col = 0
    for i, (num, title, page, color) in enumerate(sections):
        row = i % 4; col = i // 4
        x = cx + col*(cw+Inches(0.5))
        y = cy + row*(ch+Inches(0.12))
        rect(s, x, y, cw, ch, C_PANEL, color, Pt(1))
        badge(s, x+Inches(0.08), y+Inches(0.15), Inches(0.4), Inches(0.38),
              num, color, C_BG, sz=14)
        tb(s, title, x+Inches(0.6), y+Inches(0.08), cw-Inches(0.7), Inches(0.4),
           sz=12, bold=True, color=C_WHITE)
        tb(s, page, x+Inches(0.6), y+Inches(0.42), cw-Inches(0.7), Inches(0.25),
           sz=10, color=C_SUBTEXT)

    img_placeholder(s, Inches(12.0), Inches(1.3), Inches(1.1), Inches(5.9),
                    '[탭 바 클로즈업]')


def s03_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "제품 소개", "Product Overview")

    tb(s, "Universal Device Studio(UDS)는 FT232H USB 인터페이스를 통해 다양한 하드웨어 디바이스를 통합 관리하는 PC 기반 테스트 툴입니다.",
       Inches(0.3), Inches(1.15), Inches(7.8), Inches(0.55), sz=12, color=C_SUBTEXT)

    feats = [
        "I2C 버스 스캔 및 ACK/NACK 테스트",
        "SPI Probe + Device ID 검증",
        "JTAG TAP 상태 머신 다이어그램 시각화",
        "UART 시리얼 터미널 (COM 포트 자동 감지)",
        "GPIO Bit-Bang 폴링 및 레벨 제어",
        "INA228 전압/전류/전력 실시간 모니터링",
        "PI6CG18201 클럭 제너레이터 레지스터 제어",
        "ADS1018 4채널 12-bit ADC 시각화",
    ]
    card(s, Inches(0.3), Inches(1.72), Inches(5.5), Inches(3.5),
         "핵심 기능", feats, C_ACCENT)

    tb(s, "지원 모듈 (4개)", Inches(0.3), Inches(5.3), Inches(5.5), Inches(0.38),
       sz=12, bold=True, color=C_WHITE)
    mx = Inches(0.3)
    for name, color in MODULE_COLORS.items():
        badge(s, mx, Inches(5.7), Inches(2.5), Inches(0.45), name, color, C_BG)
        mx += Inches(2.65)

    img_placeholder(s, Inches(6.1), Inches(1.25), Inches(6.9), Inches(3.5),
                    '[다이어그램: PC - USB - FT232H/FT2232H/FT4232H - Target 연결 블록도]')
    img_placeholder(s, Inches(6.1), Inches(4.9), Inches(6.9), Inches(2.3),
                    '[스크린샷: 연결 완료 상태 - 초록색 LED 및 연결 정보]')


def s04_sysreq(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "시스템 요구사항 및 지원 하드웨어",
              "System Requirements & Supported Hardware")

    panels = [
        ("소프트웨어 요구사항", C_ACCENT, [
            "OS: Windows 10/11 64-bit",
            "Python 3.9+ 권장",
            "PySide6 6.7+, pyqtgraph 0.13+",
            "ftd2xx 1.3+, pyserial",
            "FTDI D2XX Driver 필수 설치",
            "pip install PySide6 pyqtgraph ftd2xx pyserial",
        ]),
        ("지원 FTDI 칩셋", C_CYAN, [
            "FT232H  —  1채널 MPSSE",
            "FT2232H  —  2채널 (A/B: MPSSE)",
            "FT4232H  —  4채널 (A/B: MPSSE, C/D: GPIO)",
            "자동 칩 감지 및 채널 구성 자동 적용",
            "C/D 채널: I2C/SPI/JTAG 비활성화 (GPIO only)",
        ]),
        ("지원 디바이스 모듈", C_GREEN, [
            "FTDI Verifier  —  프로토콜 검증 도구",
            "INA228  —  전력 측정 (I2C)",
            "PI6CG18201  —  클럭 제너레이터 (SMBus)",
            "ADS1018  —  12-bit ADC (SPI)",
            "플러그인 구조: 모듈 추가/제거 용이",
        ]),
    ]
    cx = Inches(0.3)
    for title, color, items in panels:
        card(s, cx, Inches(1.25), Inches(4.2), Inches(4.5), title, items,
             accent=color, body_sz=10.5)
        cx += Inches(4.35)

    img_placeholder(s, Inches(0.3), Inches(5.9), Inches(12.7), Inches(1.35),
                    '[사진/이미지: FT232H, FT2232H, FT4232H 보드 3종 나란히]')


def s05_install(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "설치 및 초기 설정", "Installation & Initial Launch")

    steps = [
        ("1", "드라이버 설치", C_ACCENT, [
            "FTDI 공식 사이트에서 D2XX Driver 다운로드",
            "설치 완료 후 장치 관리자 확인",
            "'USB Serial Converter' 인식 확인",
        ]),
        ("2", "패키지 설치", C_CYAN, [
            "pip install PySide6 pyqtgraph ftd2xx pyserial",
            "복수 Python 환경: where python 확인",
            "python -m pip install ... 형식 권장",
        ]),
        ("3", "앱 실행", C_GREEN, [
            "python main.py  (프로젝트 루트 기준)",
            "또는 배포 버전: DeviceManagementHub.exe",
            "콘솔 오류 없이 4개 탭 자동 로드 확인",
        ]),
        ("4", "초기 화면 확인", C_ORANGE, [
            "4개 모듈 탭 표시 확인",
            "상단 연결 툴바 활성화 확인",
            "Scan 버튼 클릭 → 장치 탐색",
        ]),
    ]
    cx = Inches(0.3)
    for num, title, color, items in steps:
        badge(s, cx+Inches(1.3), Inches(1.25), Inches(0.5), Inches(0.5),
              num, color, C_BG, sz=16)
        rect(s, cx, Inches(1.78), Inches(3.1), Inches(3.0), C_PANEL, color, Pt(1))
        rect(s, cx, Inches(1.78), Inches(0.07), Inches(3.0), color)
        tb(s, title, cx+Inches(0.15), Inches(1.84), Inches(2.85), Inches(0.35),
           sz=13, bold=True, color=color)
        rect(s, cx+Inches(0.07), Inches(2.22), Inches(2.95), Pt(0.8), color)
        y2 = Inches(2.3)
        for item in items:
            tb(s, '  > '+item, cx+Inches(0.15), y2, Inches(2.85), Inches(0.35),
               sz=10.5, color=C_TEXT)
            y2 += Inches(0.35)
        if num != "4":
            tb(s, '>>', cx+Inches(3.12), Inches(2.8), Inches(0.3), Inches(0.4),
               sz=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        cx += Inches(3.25)

    img_placeholder(s, Inches(0.3), Inches(5.0), Inches(6.3), Inches(2.2),
                    '[스크린샷: 앱 최초 실행 - 디바이스 미연결, 상태 LED 회색]')
    img_placeholder(s, Inches(6.8), Inches(5.0), Inches(6.2), Inches(2.2),
                    '[스크린샷: 장치 관리자 - FTDI D2XX 드라이버 인식 화면]')


def s06_main_interface(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "메인 인터페이스 개요", "Main Interface Overview")

    img_placeholder(s, Inches(0.3), Inches(1.2), Inches(8.5), Inches(5.6),
                    '[스크린샷: 메인 창 전체 - 번호 어노테이션 오버레이 적용\n(다크 테마, 연결 상태, FTDI Verifier 활성)]')

    callouts = [
        ("1", "Custom Title Bar",   "UDS 로고 | 테마 토글 | 최소화/최대화/닫기",  C_ACCENT),
        ("2", "Connection Toolbar", "Scan | 장치 콤보 | CH 버튼 | LED | Connect", C_CYAN),
        ("3", "Module Tab Bar",     "4개 모듈 탭 자동 로드 (플러그인 구조)",        C_GREEN),
        ("4", "Module Content",     "현재 활성 모듈의 기능 패널 (좌/우 분할)",      C_ORANGE),
        ("5", "Status Bar",         "하단: 타임스탬프 포함 상태 메시지 표시",       C_PURPLE),
    ]
    cy = Inches(1.2)
    for num, title, desc, color in callouts:
        badge(s, Inches(9.1), cy, Inches(0.42), Inches(0.38), num, color, C_BG, sz=12)
        tb(s, title, Inches(9.65), cy, Inches(3.4), Inches(0.35),
           sz=11, bold=True, color=color)
        tb(s, desc, Inches(9.65), cy+Inches(0.34), Inches(3.4), Inches(0.3),
           sz=9.5, color=C_SUBTEXT)
        cy += Inches(1.1)


def s07_connection_toolbar(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "연결 툴바 상세", "Connection Toolbar Walkthrough")

    img_placeholder(s, Inches(0.3), Inches(1.2), Inches(12.7), Inches(1.1),
                    '[스크린샷: 연결 툴바 클로즈업 - FT2232H 연결, Channel B 선택됨, 초록 LED]')

    items = [
        ("Scan (R)", C_ACCENT,
         "FTDI 장치 전체 열거. 멀티채널은 채널별 분리 표시. 연결 중에는 비활성화."),
        ("Device Combo", C_CYAN,
         "감지된 장치 목록 (시리얼번호 + 모델명 + 채널수). 4-tuple 데이터 내부 저장."),
        ("CH A/B/C/D", C_GREEN,
         "멀티채널 장치에서 사용 채널 선택. 단채널 자동 선택. FT4232H C/D = MPSSE 불가 표시."),
        ("Status LED", C_ORANGE,
         "회색=미연결 | 초록=연결됨 | 빨간=오류/연결실패"),
        ("Connect", C_PURPLE,
         "선택 장치+채널로 MPSSE 세션 오픈. 연결 후 'Disconnect' 버튼으로 전환."),
        ("Connection Info", C_YELLOW,
         "연결 완료 시 'Connected: SN=xxxx, CH=A' 형식으로 상세 정보 표시."),
    ]
    cx, cy = Inches(0.3), Inches(2.5)
    cw, ch = Inches(4.15), Inches(1.45)
    for i, (title, color, desc) in enumerate(items):
        col, row = i % 3, i // 3
        x = cx + col*(cw+Inches(0.1))
        y = cy + row*(ch+Inches(0.1))
        rect(s, x, y, cw, ch, C_PANEL, color, Pt(1))
        rect(s, x, y, Inches(0.07), ch, color)
        tb(s, title, x+Inches(0.12), y+Inches(0.07), cw-Inches(0.2), Inches(0.38),
           sz=12, bold=True, color=color)
        tb(s, desc, x+Inches(0.12), y+Inches(0.45), cw-Inches(0.2), Inches(0.9),
           sz=10, color=C_TEXT)

    img_placeholder(s, Inches(9.15), Inches(5.4), Inches(4.0), Inches(1.85),
                    '[스크린샷: 콤보박스 드롭다운 - 복수 장치 목록]')


def s08_verifier_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "FTDI Verifier  —  모듈 개요",
              "Hardware Verification Tool | I2C / SPI / JTAG / UART / GPIO")

    img_placeholder(s, Inches(0.3), Inches(1.2), Inches(7.5), Inches(5.6),
                    '[스크린샷: FTDI Verifier 탭 전체\n(I2C 모드, FT232H 핀아웃 우측 표시)]')

    feats = [
        "실시간 핀아웃 시각화 (hover/click 인터랙션)",
        "I2C Bus Scan (0x08~0x77, 이력 테이블)",
        "SPI Probe + Device ID 검증",
        "JTAG TAP 상태 머신 다이어그램",
        "UART 터미널 (pyserial 기반 COM 포트)",
        "GPIO Bit-Bang 폴링 (50ms~10000ms)",
        "채널 전환 시 칩 스펙 자동 재적용",
        "FT4232H C/D: MPSSE 불가 구간 자동 비활성화",
    ]
    card(s, Inches(8.1), Inches(1.25), Inches(5.0), Inches(3.5),
         "주요 기능", feats, C_ACCENT)

    proto = [("I2C",  C_ACCENT), ("SPI", C_ORANGE), ("JTAG", C_PURPLE),
             ("UART", C_GREEN),  ("GPIO", C_YELLOW)]
    px = Inches(8.1)
    for name, color in proto:
        badge(s, px, Inches(4.9), Inches(0.92), Inches(0.42), name, color, C_BG, sz=11)
        px += Inches(1.01)

    img_placeholder(s, Inches(8.1), Inches(5.45), Inches(5.0), Inches(1.35),
                    '[스크린샷: Device Info 패널 - Chip: FT2232H, Channel: A]')


def s09_i2c(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "FTDI Verifier  —  I2C 탭", "I2C Bus Scan & ACK Test")

    two_col_header(s, "I2C 컨트롤 패널", "스캔 결과 테이블", C_ACCENT, C_CYAN)

    left_items = [
        "I2C Bus Scan 버튼: 0x08~0x77 전체 스캔",
        "스캔 소요 시간 자동 측정 및 표시",
        "Scanned Address 콤보: 스캔 결과 자동 채워짐",
        "콤보박스 직접 입력: 16진수 주소 수동 입력 가능",
        "ACK TEST 버튼: 선택 주소로 I2C 프레임 전송",
        "ACK LED: 초록=ACK / 빨간=NACK / 회색=N/A",
    ]
    bullet_list(s, left_items, Inches(0.3), Inches(1.7), Inches(6.1), sz=11)

    # Right: table mockup
    cols = ["Address", "Status"]
    widths = [Inches(2.8), Inches(3.5)]
    rx = Inches(6.55)
    ry = Inches(1.7)
    rect(s, rx, ry, Inches(6.5), Inches(0.38), C_ACCENT)
    for i, (col, w) in enumerate(zip(cols, widths)):
        tb(s, col, rx, ry, w, Inches(0.38), sz=11, bold=True,
           color=C_BG, align=PP_ALIGN.CENTER)
        rx += w
    rows = [("0x40 (INA228)", "ACK OK"), ("0x48 (ADS1115)", "ACK OK"),
            ("0x50 (PI6CG)", "ACK OK"), ("0x76 (BME280)", "ACK OK")]
    rx = Inches(6.55)
    for i, (addr, stat) in enumerate(rows):
        ry2 = ry + Inches(0.38) + i*Inches(0.36)
        bg = C_PANEL2 if i % 2 == 0 else C_PANEL
        rect(s, rx, ry2, Inches(6.5), Inches(0.36), bg)
        tb(s, addr, rx+Inches(0.1), ry2, Inches(2.7), Inches(0.36),
           sz=10.5, color=C_TEXT)
        badge(s, rx+Inches(3.1), ry2+Inches(0.03), Inches(1.4), Inches(0.3),
              stat, C_GREEN, C_BG, sz=9)

    img_placeholder(s, Inches(6.55), Inches(3.35), Inches(6.5), Inches(1.5),
                    '[스크린샷: I2C 스캔 이력 테이블 - 타임스탬프, 발견 주소, 소요시간]')
    img_placeholder(s, Inches(0.3), Inches(4.95), Inches(6.1), Inches(2.25),
                    '[스크린샷: ACK TEST - ACK LED 초록색 점등 클로즈업]')

    # Notes
    tb(s, "주의: FT4232H CH C/D 선택 시 I2C/SPI/JTAG 버튼 자동 비활성화",
       Inches(6.55), Inches(4.95), Inches(6.5), Inches(0.38),
       sz=10, color=C_YELLOW, italic=True)
    tb(s, "슬레이브 주소는 7비트 형식으로 입력 (FtdiManager 내부에서 <<1 시프트 자동 처리)",
       Inches(6.55), Inches(5.35), Inches(6.5), Inches(0.38),
       sz=10, color=C_SUBTEXT, italic=True)


def s10_spi(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "FTDI Verifier  —  SPI 탭", "SPI Probe & Device ID Verification")

    section_header(s, "SPI 설정", Inches(1.15), C_ORANGE)
    modes = [
        ("Mode 0", "CPOL=0 CPHA=0", C_GREEN),
        ("Mode 1", "CPOL=0 CPHA=1", C_CYAN),
        ("Mode 2", "CPOL=1 CPHA=0", C_ORANGE),
        ("Mode 3", "CPOL=1 CPHA=1", C_PURPLE),
    ]
    mx = Inches(0.3)
    for mode, desc, color in modes:
        rect(s, mx, Inches(1.58), Inches(2.95), Inches(0.75), C_PANEL, color, Pt(1))
        tb(s, mode, mx+Inches(0.1), Inches(1.62), Inches(2.75), Inches(0.35),
           sz=13, bold=True, color=color)
        tb(s, desc, mx+Inches(0.1), Inches(1.97), Inches(2.75), Inches(0.3),
           sz=10, color=C_SUBTEXT)
        mx += Inches(3.07)

    section_header(s, "SPI 클럭 속도", Inches(2.45), C_ORANGE)
    freqs = ["100 kHz", "400 kHz", "1 MHz", "2 MHz", "4 MHz", "6 MHz", "10 MHz", "12 MHz"]
    fx = Inches(0.3)
    for f in freqs:
        badge(s, fx, Inches(2.88), Inches(1.5), Inches(0.36), f, C_PANEL2, C_TEXT, sz=10)
        fx += Inches(1.55)

    img_placeholder(s, Inches(0.3), Inches(3.35), Inches(6.1), Inches(1.5),
                    '[스크린샷: SPI 타이밍 다이어그램 - Mode 0 CLK/MOSI/MISO 파형 뷰]')

    section_header(s, "SPI Probe  &  Device ID 검증", Inches(5.0), C_ORANGE)
    probe_items = [
        "SPI Probe 버튼: CS 활성화 + 더미 프레임 전송 → 응답 유무 확인",
        "Register 필드: 읽을 레지스터 주소 입력 (예: 0x9F JEDEC ID)",
        "Bytes 콤보: 읽을 바이트 수 선택 (1~4 bytes)",
        "Expected 필드: 예상 값 입력 (불일치 시 경고 표시)",
        "Read ID 버튼 → 읽기 결과 라벨 표시",
    ]
    bullet_list(s, probe_items, Inches(0.3), Inches(5.45), Inches(6.1), sz=11)

    img_placeholder(s, Inches(6.55), Inches(3.35), Inches(6.5), Inches(3.85),
                    '[스크린샷: SPI Device ID 검증 - 0x9F 읽기 결과 "0xEF4018 MATCH" 표시]')


def s11_jtag(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "FTDI Verifier  —  JTAG 탭", "TAP State Machine Diagram & Pin Control")

    img_placeholder(s, Inches(3.7), Inches(1.2), Inches(9.3), Inches(5.6),
                    '[스크린샷: JTAG 모드 - TAP 다이어그램 + Pin Status 사이드바 전체\n(현재 상태 하이라이트, TCK/TMS/TDI/TDO/RST LED 표시)]')

    left_items = [
        "TAP 상태 머신 다이어그램",
        "16개 상태 노드 + 전이 화살표",
        "4K HiDPI 고해상도 PNG 렌더링",
        "현재 상태 하이라이트 표시",
        "",
        "Pin Status 사이드바",
        "TCK / TMS / TDI / TDO / RST",
        "각 핀 LED 점등 + HIGH/LOW 텍스트",
        "",
        "Mapping Management",
        "CSV Import/Export 지원",
        "Signal-to-Pin 매핑 테이블 (5열)",
        "",
        "TDO Logger",
        "수신 데이터 실시간 로그",
    ]
    y = Inches(1.2)
    for item in left_items:
        if item == "":
            y += Inches(0.12)
            continue
        is_title = item in ["TAP 상태 머신 다이어그램", "Pin Status 사이드바",
                             "Mapping Management", "TDO Logger"]
        tb(s, item, Inches(0.3), y, Inches(3.2), Inches(0.33),
           sz=11, bold=is_title, color=C_PURPLE if is_title else C_TEXT)
        y += Inches(0.33)

    img_placeholder(s, Inches(0.3), Inches(5.5), Inches(3.2), Inches(1.3),
                    '[스크린샷: CSV 가져온 후 핀 매핑 테이블]')


def s12_uart_gpio(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "FTDI Verifier  —  UART  &  GPIO 탭",
              "Serial Terminal & GPIO Bit-Bang Control")

    two_col_header(s, "UART 탭", "GPIO 탭", C_GREEN, C_YELLOW)

    uart_items = [
        "COM 포트 자동 감지 (pyserial) + 새로고침",
        "Baudrate: 9600 ~ 921600 (8단계)",
        "Data bits: 7/8 | Parity: None/Even/Odd/Mark/Space",
        "Stop bits: 1 / 1.5 / 2",
        "Flow control: None / RTS-CTS / XON-XOFF",
        "OPEN/CLOSE 버튼 (전체 너비 토글)",
        "콘솔: RX Format(ASCII/HEX), AutoScroll, Timestamp",
        "3열 테이블: Time / Dir(TX/RX) / Data",
        "Send 입력 + EOL 콤보 (No/CR/LF/CRLF) + SEND",
        "[주의] 탭 이탈 시 COM 포트 자동 닫힘 + MPSSE 복원",
    ]
    bullet_list(s, uart_items, Inches(0.3), Inches(1.7), Inches(6.1), sz=10.5)

    gpio_items = [
        "Polling Interval: 50ms ~ 10000ms (SpinBox)",
        "Start/Stop Polling 토글 버튼",
        "GPIO Backend: BITBANG 모드 표시",
        "핀별 토글 버튼: 'GPIO: HIGH' / 'GPIO: LOW' 레이블",
        "D0~D7 8핀 Bit-Bang 그리드",
        "GPIO 상태 테이블 (5열): Pin / Name / Mode / Dir / Level",
        "[주의] 크로스 채널 동작 시 경고 표시",
    ]
    bullet_list(s, gpio_items, Inches(6.6), Inches(1.7), Inches(6.5), sz=10.5)

    img_placeholder(s, Inches(0.3), Inches(5.1), Inches(6.1), Inches(2.15),
                    '[스크린샷: UART 탭 - COM3 열린 상태, 데이터 수신 중인 콘솔 테이블]')
    img_placeholder(s, Inches(6.6), Inches(5.1), Inches(6.5), Inches(2.15),
                    '[스크린샷: GPIO 탭 - Start Polling 활성, D0~D3 HIGH/LOW 혼합 상태]')


def s13_ina228_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "INA228 전력 모니터  —  모듈 개요",
              "Real-time Voltage / Current / Power / Temperature Monitor")

    img_placeholder(s, Inches(4.5), Inches(1.2), Inches(8.6), Inches(5.6),
                    '[스크린샷: INA228 탭 전체\n(모니터링 중, 이중 차트에 실시간 데이터 표시)]')

    panels = [
        ("주소 스캔", C_CYAN, [
            "Scan Addresses 버튼",
            "범위: 0x40 ~ 0x4F",
            "감지 콤보박스 자동 채워짐",
        ]),
        ("ADC 설정", C_ACCENT, [
            "ADC Range: 80mV / 40mV",
            "AVG Count: 1 ~ 1024",
            "변환 시간: 50us ~ 4.12ms",
            "Shunt 저항값 입력 (mOhm)",
        ]),
        ("GPIO Hold", C_ORANGE, [
            "D4~D7 토글 버튼",
            "HIGH/OFF 레벨 잠금",
            "LED + Progress Bar 표시",
        ]),
    ]
    cy = Inches(1.25)
    for title, color, items in panels:
        card(s, Inches(0.3), cy, Inches(4.0), len(items)*Inches(0.35)+Inches(0.85),
             title, items, color, body_sz=10.5)
        cy += len(items)*Inches(0.35)+Inches(0.85)+Inches(0.15)


def s14_ina228_chart(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "INA228  —  실시간 차트 & 메트릭",
              "Dual pyqtgraph Chart | Voltage & Current Live Monitoring")

    img_placeholder(s, Inches(0.3), Inches(1.2), Inches(7.8), Inches(3.4),
                    '[스크린샷: pyqtgraph 이중 차트 클로즈업\n(상단: 전압 3.3V 안정적, 하단: 전류 파형)]')

    metrics = [
        ("VBUS", "버스 전압", "V", C_ACCENT),
        ("VSHUNT", "션트 전압", "mV", C_CYAN),
        ("Current", "전류", "A", C_ORANGE),
        ("Power", "소비 전력", "mW", C_GREEN),
    ]
    mx = Inches(0.3)
    for name, desc, unit, color in metrics:
        rect(s, mx, Inches(4.75), Inches(1.87), Inches(1.1), C_PANEL, color, Pt(1))
        tb(s, name, mx+Inches(0.08), Inches(4.82), Inches(1.7), Inches(0.38),
           sz=13, bold=True, color=color)
        tb(s, desc+' ('+unit+')', mx+Inches(0.08), Inches(5.18), Inches(1.7), Inches(0.3),
           sz=9, color=C_SUBTEXT)
        mx += Inches(1.97)

    badge(s, Inches(8.05), Inches(4.75), Inches(1.6), Inches(1.1),
          "Die Temp\n(deg C)", C_PURPLE, C_WHITE, sz=11)

    tech_items = [
        "슬라이딩 윈도우: 기본 60초, 최대 2000 포인트",
        "DIAG_ALRT[bit0] 폴링으로 변환 완료 감지 (CNVRF)",
        "0-스파이크 필터: 이전값 비영에서 새값이 0이면 자동 스킵",
        "VSHUNT 변환: 3바이트 >> 4 (20-bit), Range0=312.5nV LSB",
        "VBUS 변환: 3바이트 >> 4 (20-bit), LSB=195.3125uV",
        "DIETEMP 변환: 2바이트, LSB=7.8125 m-degC",
        "Window 설정 (10~600s): 데이터 보존 구간 조절",
    ]
    card(s, Inches(8.2), Inches(1.2), Inches(4.9), Inches(4.55),
         "기술 세부사항", tech_items, C_CYAN, body_sz=10)


def s15_ina228_regmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "INA228  —  레지스터 맵", "INA228 Register Reference")

    img_placeholder(s, Inches(7.5), Inches(1.2), Inches(5.6), Inches(5.6),
                    '[스크린샷: INA228 레지스터 맵 테이블\n(decoded 값과 hex 동시 표시)]')

    headers = ["레지스터", "주소", "크기", "설명"]
    widths = [Inches(2.3), Inches(0.9), Inches(0.7), Inches(3.3)]
    regs = [
        ("CONFIG",        "0x00", "2B", "ADC 리셋, 변환 지연, 온도 보정 활성화"),
        ("ADC_CONFIG",    "0x01", "2B", "MODE, VBUSCT, VSHCT, SHCT, AVG 설정"),
        ("SHUNT_CAL",     "0x02", "2B", "전류 계산 스케일 팩터 (CURRENT_LSB 기반)"),
        ("VSHUNT",        "0x04", "3B", "션트 차동 전압 (20-bit, 부호 있음)"),
        ("VBUS",          "0x05", "3B", "버스 전압 (20-bit, 부호 없음)"),
        ("DIETEMP",       "0x06", "2B", "다이 온도 (16-bit, 12-bit 유효)"),
        ("CURRENT",       "0x07", "3B", "계산된 전류값 (SHUNT_CAL 적용)"),
        ("POWER",         "0x08", "3B", "계산된 전력값"),
        ("DIAG_ALRT",     "0x0B", "2B", "변환 완료(CNVRF) + 알람 플래그"),
        ("MANUFACTURER",  "0x3E", "2B", "제조사 ID (0x5449 = Texas Instruments)"),
        ("DEVICE_ID",     "0x3F", "2B", "디바이스 ID (0x2281)"),
    ]
    ry = Inches(1.25)
    rx0 = Inches(0.3)
    # Header
    rect(s, rx0, ry, sum(widths), Inches(0.38), C_ACCENT)
    rx = rx0
    for h, w in zip(headers, widths):
        tb(s, h, rx, ry, w, Inches(0.38), sz=11, bold=True,
           color=C_BG, align=PP_ALIGN.CENTER)
        rx += w
    ry += Inches(0.38)
    for i, (name, addr, size, desc) in enumerate(regs):
        bg = C_PANEL2 if i % 2 == 0 else C_PANEL
        rect(s, rx0, ry, sum(widths), Inches(0.36), bg)
        rx = rx0
        for val, w in zip([name, addr, size, desc], widths):
            tb(s, val, rx+Inches(0.05), ry+Inches(0.03), w-Inches(0.1), Inches(0.3),
               sz=10, color=C_CYAN if i % 2 == 0 else C_TEXT)
            rx += w
        ry += Inches(0.36)


def s16_pi6cg_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "PI6CG18201  —  클럭 제너레이터",
              "Programmable Differential Clock Generator Control")

    two_col_header(s, "제어 패널 (Control Center)", "클럭 파형 시각화", C_GREEN, C_CYAN)

    ctrl_items = [
        "Slave Address: 0x68 (SADR=L) / 0x6A (SADR=H)",
        "Advanced Mode: 바이트/비트 매핑 힌트 표시",
        "",
        "Output Enable (OE):",
        "  Q0 (Q0+/Q0-) 체크박스 - 청록색",
        "  Q1 (Q1+/Q1-) 체크박스 - 분홍색",
        "",
        "Output Amplitude: 0.6V / 0.7V / 0.8V / 0.9V",
        "",
        "Spread Spectrum (SS):",
        "  Auto (HW 핀 SS0/SS1 추종)",
        "  Off / Weak (-0.25%) / Strong (-0.5%)",
        "  SS Readback 뱃지 (현재 HW 상태)",
        "",
        "Slew Rate (Coarse): Both Slow ~ Both Fast",
        "Slew Rate (Fine REF): Slow/Medium/Fast/Very Fast",
        "Combined Level: 'Lv.8/15' 형식 표시",
        "",
        "Live Mode: ON 시 설정 변경 즉시 I2C 전송",
        "Read / Write 버튼 (수동 트리거)",
    ]
    y = Inches(1.7)
    for item in ctrl_items:
        if item == "":
            y += Inches(0.1); continue
        is_section = item.endswith(':') and not item.startswith(' ')
        tb(s, item, Inches(0.3), y, Inches(6.1), Inches(0.3),
           sz=10.5 if not is_section else 11,
           bold=is_section, color=C_GREEN if is_section else C_TEXT)
        y += Inches(0.3)

    img_placeholder(s, Inches(6.55), Inches(1.7), Inches(6.5), Inches(2.8),
                    '[스크린샷: 클럭 파형 시각화 - Q0/Q1 차동 파형, 진폭/슬루율 반영]')
    img_placeholder(s, Inches(6.55), Inches(4.65), Inches(6.5), Inches(2.1),
                    '[스크린샷: 제어 패널 - Advanced Mode ON, 힌트 텍스트 표시 중]')


def s17_pi6cg_regmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "PI6CG18201  —  레지스터 맵 & Live Mode",
              "Register Bit-Field Reference | SMBus Protocol")

    img_placeholder(s, Inches(7.5), Inches(1.2), Inches(5.6), Inches(3.0),
                    '[스크린샷: 레지스터 맵 테이블 - Byte0~7 hex + decoded 비트필드]')
    img_placeholder(s, Inches(7.5), Inches(4.35), Inches(5.6), Inches(2.9),
                    '[스크린샷: Live Mode ON 상태 - I2C 로그에 Write 기록 실시간]')

    headers = ["BitField", "Byte", "Bits", "설명"]
    widths = [Inches(1.8), Inches(0.7), Inches(0.8), Inches(3.9)]
    fields = [
        ("OE_Q0",     "0", "[1]",    "Q0 출력 인에이블"),
        ("OE_Q1",     "0", "[2]",    "Q1 출력 인에이블"),
        ("AMPLITUDE", "1", "[1:0]",  "출력 전압 레벨 (0.6/0.7/0.8/0.9V)"),
        ("SS_SW_CTRL","1", "[5]",    "Spread Spectrum 소프트웨어 제어 활성"),
        ("SS_MODE",   "1", "[4:3]",  "SS 모드 (Off/Weak/Strong)"),
        ("SLEW_Q0",   "2", "[1]",    "Q0 슬루율 (0=Slow, 1=Fast)"),
        ("SLEW_Q1",   "2", "[2]",    "Q1 슬루율"),
        ("REF_SLEW",  "3", "[7:6]",  "REF 슬루율 (Slow~Very Fast)"),
    ]
    ry = Inches(1.25)
    rect(s, Inches(0.3), ry, sum(widths), Inches(0.38), C_GREEN)
    rx = Inches(0.3)
    for h, w in zip(headers, widths):
        tb(s, h, rx, ry, w, Inches(0.38), sz=11, bold=True,
           color=C_BG, align=PP_ALIGN.CENTER); rx += w
    ry += Inches(0.38)
    for i, (name, byte, bits, desc) in enumerate(fields):
        bg = C_PANEL2 if i % 2 == 0 else C_PANEL
        rect(s, Inches(0.3), ry, sum(widths), Inches(0.36), bg)
        for val, w, lx in zip([name, byte, bits, desc], widths,
                              [Inches(0.3), Inches(2.1), Inches(2.8), Inches(3.6)]):
            tb(s, val, lx+Inches(0.05), ry+Inches(0.03), w-Inches(0.1), Inches(0.3),
               sz=10, color=C_GREEN if i % 2 == 0 else C_TEXT)
        ry += Inches(0.36)

    # SMBus protocol note
    tb(s, "SMBus 프로토콜:", Inches(0.3), Inches(4.45), Inches(7.0), Inches(0.35),
       sz=12, bold=True, color=C_GREEN)
    smb = [
        "Write: [SlaveAddr W] [0x00] [Byte0 ~ Byte7] (8 bytes block write)",
        "Read:  [SlaveAddr W] [0x00] [SlaveAddr R] [8 bytes block read]",
        "Live Mode OFF: 변경 누적 -> Write 버튼으로 일괄 전송",
        "Live Mode ON:  설정 변경 즉시 SMBus Block Write 실행",
    ]
    bullet_list(s, smb, Inches(0.3), Inches(4.85), Inches(7.0), sz=10.5)


def s18_ads1018_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "ADS1018 ADC 모니터  —  모듈 개요",
              "4-Channel 12-bit SPI ADC Real-time Monitor")

    img_placeholder(s, Inches(4.5), Inches(1.2), Inches(8.6), Inches(5.6),
                    '[스크린샷: ADS1018 탭 전체\n(4채널 실시간 ADC 파형, CH0=전압/CH1=전류 모드)]')

    adc_settings = [
        "PGA (게인): +/-6.144V ~ +/-0.256V (6단계)",
        "Data Rate: 128 ~ 3300 SPS (7단계)",
        "Mode: Continuous / Single-shot",
        "Sensor: ADC / 내장 온도 센서",
        "CS Pin: ADBUS3 ~ ADBUS7 선택",
        "Pull-up 체크박스",
        "Window: 10 ~ 600초 (슬라이딩 버퍼)",
        "Auto Range: Y축 자동 스케일",
    ]
    card(s, Inches(0.3), Inches(1.25), Inches(3.95), Inches(3.6),
         "ADC 설정", adc_settings, C_ORANGE, body_sz=10)

    ch_items = [
        "CH0 ~ CH3 각각 독립 설정",
        "V 모드: 전압 측정 (V)",
        "I 모드: 전류 측정 (mA)",
        "I 모드 시: R(션트저항) + G(게인) 입력",
        "예: R=0.02Ohm, G=100 -> 1mV = 0.5mA",
    ]
    card(s, Inches(0.3), Inches(5.0), Inches(3.95), Inches(2.25),
         "채널 설정 (CH0~CH3)", ch_items, C_CYAN, body_sz=10)


def s19_ads1018_tech(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "ADS1018  —  기술 상세", "Config Register & Channel Conversion Reference")

    two_col_header(s, "Config Register 구조", "전압/전류 변환 공식", C_ORANGE, C_CYAN)

    headers = ["Bits", "Field", "설명"]
    widths = [Inches(1.0), Inches(1.4), Inches(3.8)]
    cfg_regs = [
        ("[15]",   "OS",     "Single-shot 변환 시작"),
        ("[14:12]","MUX",    "입력 채널/차동 쌍 선택 (AIN0~3)"),
        ("[11:9]", "PGA",    "프로그래머블 게인 설정"),
        ("[8]",    "MODE",   "Continuous / Single-shot"),
        ("[7:5]",  "DR",     "데이터 레이트 (128~3300 SPS)"),
        ("[4]",    "TS_MODE","ADC (0) / 온도 센서 (1)"),
        ("[3]",    "PULL_UP","DOUT 풀업 활성화"),
        ("[2:1]",  "NOP",    "No-op (필수값: 0b01)"),
    ]
    ry = Inches(1.7)
    rect(s, Inches(0.3), ry, sum(widths), Inches(0.35), C_ORANGE)
    rx = Inches(0.3)
    for h, w in zip(headers, widths):
        tb(s, h, rx, ry, w, Inches(0.35), sz=11, bold=True,
           color=C_BG, align=PP_ALIGN.CENTER); rx += w
    ry += Inches(0.35)
    for i, (bits, field, desc) in enumerate(cfg_regs):
        bg = C_PANEL2 if i % 2 == 0 else C_PANEL
        rect(s, Inches(0.3), ry, sum(widths), Inches(0.33), bg)
        for val, w, lx in zip([bits, field, desc], widths,
                              [Inches(0.3), Inches(1.3), Inches(2.7)]):
            tb(s, val, lx+Inches(0.05), ry+Inches(0.02), w-Inches(0.1), Inches(0.29),
               sz=10, color=C_ORANGE if i % 2 == 0 else C_TEXT)
        ry += Inches(0.33)

    # Right side
    conv_items = [
        "전압 모드: V = raw_12bit * LSB_mV / 1000",
        "PGA +/-2.048V 기준: LSB = 1.0 mV",
        "PGA +/-0.256V 기준: LSB = 0.125 mV",
        "",
        "전류 모드: I = V / (R_shunt * Gain)",
        "예) R=0.02Ohm, G=100 -> 1mV = 0.5 mA",
        "채널별 독립 V/I 혼합 사용 가능",
        "",
        "SPI 프로토콜: CPOL=0, CPHA=1 (Mode 1)",
        "Config Register: 16-bit 쓰기/읽기",
        "Throughput: 4채널 = 실질 약 825 SPS/ch",
    ]
    y = Inches(1.7)
    for item in conv_items:
        if item == "":
            y += Inches(0.12); continue
        tb(s, item, Inches(6.6), y, Inches(6.4), Inches(0.33),
           sz=10.5, color=C_TEXT)
        y += Inches(0.33)

    img_placeholder(s, Inches(0.3), Inches(5.65), Inches(12.7), Inches(1.6),
                    '[스크린샷: 4채널 pyqtgraph 파형 - 각 채널 고유 색상, 측정값 메트릭 카드 표시]')


def s20_troubleshoot(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "트러블슈팅 & 주의사항", "Troubleshooting & Important Notes")

    groups = [
        ("연결 문제", C_RED, [
            ("Scan 후 목록 비어 있음",
             "D2XX 드라이버 미설치 또는 VCP 드라이버 충돌. 드라이버 재설치 후 장치 관리자 확인."),
            ("모듈 탭 미표시",
             "의존성 패키지 누락 (ModuleNotFoundError). 콘솔 로그 확인 후 pip install 재실행."),
            ("FT4232H에서 I2C/SPI 미동작",
             "C/D 채널은 MPSSE 미지원. A 또는 B 채널로 변경 후 재연결."),
        ]),
        ("통신 주의사항", C_YELLOW, [
            ("UART 탭 이탈 시 포트 닫힘",
             "의도된 동작: 다른 탭 전환 시 COM 포트 자동 닫힘 + FTDI MPSSE 복원."),
            ("I2C 주소 입력 형식",
             "7비트 주소 입력 (0x40 등). 내부에서 <<1 시프트 자동 처리. 8비트 입력 금지."),
            ("GPIO Hold와 I2C 충돌",
             "INA228 모니터링 중 D4~D7 Hold가 SDA/SCL에 영향 가능. Hold 핀 확인 필수."),
        ]),
        ("성능 & 데이터", C_GREEN, [
            ("INA228 0-스파이크 필터",
             "이전 유효값이 비영인데 새 값이 0이면 자동 스킵 (노이즈 제거 목적)."),
            ("pyqtgraph 버퍼 한계",
             "최대 2000 포인트 슬라이딩. Window(초) 줄이면 해상도 향상."),
            ("한글 소스 인코딩",
             "소스 수정 시 한글은 \\uXXXX 유니코드 이스케이프 사용 권장 (CP949 오염 방지)."),
        ]),
    ]
    cx = Inches(0.3)
    cw = Inches(4.2)
    for group_title, color, items in groups:
        rect(s, cx, Inches(1.2), cw, Inches(0.38), color)
        tb(s, group_title, cx, Inches(1.2), cw, Inches(0.38),
           sz=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        cy = Inches(1.62)
        for symptom, solution in items:
            rect(s, cx, cy, cw, Inches(1.35), C_PANEL, color, Pt(0.8))
            tb(s, symptom, cx+Inches(0.1), cy+Inches(0.07),
               cw-Inches(0.2), Inches(0.35), sz=11, bold=True, color=color)
            tb(s, solution, cx+Inches(0.1), cy+Inches(0.42),
               cw-Inches(0.2), Inches(0.85), sz=10, color=C_TEXT)
            cy += Inches(1.42)
        cx += cw + Inches(0.17)

    # Quick reference table
    tb(s, "Quick Reference", Inches(0.3), Inches(5.6), Inches(5), Inches(0.38),
       sz=12, bold=True, color=C_WHITE)
    qr = [
        ("장치 연결",   "Scan -> 장치 선택 -> 채널 선택 -> Connect"),
        ("I2C 스캔",    "FTDI Verifier -> I2C -> Bus Scan"),
        ("전력 모니터", "INA228 -> Scan Addresses -> Start"),
        ("클럭 설정",   "PI6CG18201 -> Read -> 수정 -> Write"),
        ("ADC 측정",    "ADS1018 -> 채널 설정 -> Start"),
    ]
    qx = Inches(0.3); qy = Inches(6.05)
    for action, steps in qr:
        rect(s, qx, qy, Inches(2.4), Inches(0.3), C_ACCENT)
        tb(s, action, qx, qy, Inches(2.4), Inches(0.3), sz=9, bold=True,
           color=C_BG, align=PP_ALIGN.CENTER)
        rect(s, qx+Inches(2.4), qy, Inches(10.2), Inches(0.3), C_PANEL2)
        tb(s, steps, qx+Inches(2.5), qy, Inches(10.0), Inches(0.3), sz=9, color=C_TEXT)
        qy += Inches(0.32)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        (s01_cover,           "Cover"),
        (s02_toc,             "Table of Contents"),
        (s03_overview,        "Product Overview"),
        (s04_sysreq,          "System Requirements"),
        (s05_install,         "Installation"),
        (s06_main_interface,  "Main Interface"),
        (s07_connection_toolbar, "Connection Toolbar"),
        (s08_verifier_overview,  "FTDI Verifier Overview"),
        (s09_i2c,             "Verifier I2C"),
        (s10_spi,             "Verifier SPI"),
        (s11_jtag,            "Verifier JTAG"),
        (s12_uart_gpio,       "Verifier UART & GPIO"),
        (s13_ina228_overview, "INA228 Overview"),
        (s14_ina228_chart,    "INA228 Chart"),
        (s15_ina228_regmap,   "INA228 Register Map"),
        (s16_pi6cg_overview,  "PI6CG18201 Overview"),
        (s17_pi6cg_regmap,    "PI6CG18201 Registers"),
        (s18_ads1018_overview,"ADS1018 Overview"),
        (s19_ads1018_tech,    "ADS1018 Technical"),
        (s20_troubleshoot,    "Troubleshooting"),
    ]

    print(f"Generating UDS User Manual PPT  ({len(slides)} slides)...")
    for i, (fn, name) in enumerate(slides, 1):
        fn(prs)
        print(f"  [{i:2d}/{len(slides)}] {name}")

    out_dir = os.path.join(os.path.dirname(__file__), "..", "docs")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "UDS_UserManual.pptx")
    prs.save(out_path)
    print(f"\nSaved: {os.path.abspath(out_path)}")
    print(f"Image placeholders: 34 locations (insert screenshots manually)")


if __name__ == "__main__":
    main()
